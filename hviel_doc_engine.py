"""
hviel_doc_engine.py
====================
THE MISSING PIECE — This is what makes your $10 API key generate files like Claude does.

HOW CLAUDE.AI WORKS (what you're paying $20/month for):
    [User asks for file] → Claude writes code → Server runs code → File created

HOW YOUR HVIEL WORKS (with this script):
    [User asks for file] → Claude API returns content (JSON) → THIS SCRIPT builds the file

INSTALL:
    pip install python-docx openpyxl python-pptx reportlab anthropic

Author: Built for Hviel / PRC Libya
"""

import os
import json
import logging
import time

_logger = logging.getLogger("hviel-doc")
from datetime import datetime

# ── File builders ──
from docx import Document as DocxDocument
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph as RLParagraph, Spacer,
    Table as RLTable, TableStyle
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT


# ═══════════════════════════════════════════════════════
# PRC BRAND CONSTANTS
# ═══════════════════════════════════════════════════════
NAVY        = "1B3A5C"
BLUE        = "2E75B6"
ZEBRA       = "F2F2F2"
LIGHT_BLUE  = "D5E8F0"
WHITE       = "FFFFFF"
DARK_GRAY   = "333333"

NAVY_RGB       = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_RGB       = RGBColor(0x2E, 0x75, 0xB6)
GRAY_RGB       = RGBColor(0x33, 0x33, 0x33)
WHITE_RGB      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_RGB = RGBColor(0x99, 0x99, 0x99)
RED_RGB        = RGBColor(0xC0, 0x39, 0x2B)


import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# PRC brand colour palette for multi-curve plots
_PRC_COLORS = ['#1e3a8a', '#E31E24', '#F59E0B', '#16a34a', '#7c3aed', '#0891b2']

# ═══════════════════════════════════════════════════════
# MAIN ENGINE CLASS
# ═══════════════════════════════════════════════════════
class HvielDocEngine:
    """
    Bridge between Claude API and actual file generation.

    Claude API returns structured JSON → this class builds .docx/.xlsx/.pptx/.pdf
    """

    def __init__(self, api_key: str = None, output_dir: str = "."):
        self.api_key    = api_key
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def _extract_json_payload(self, text: str) -> tuple[dict, str]:
        """
        Extracts the JSON payload from a __PRC_PLOT__ tag using a bracket-counting parser,
        and returns a tuple of (parsed_dict, trailing_text).
        """
        if '__PRC_PLOT__' not in text:
            return {}, text
            
        parts = text.split('__PRC_PLOT__', 1)
        leading_text = parts[0].strip()
        rest = parts[1].strip()
        
        # Clean potential markdown block prefix
        if rest.startswith('```json'):
            rest = rest[7:].strip()
        elif rest.startswith('```'):
            rest = rest[3:].strip()
            
        # Find the starting brace
        start_idx = rest.find('{')
        if start_idx == -1:
            return {}, text
            
        brace_count = 0
        in_string = False
        escape_next = False
        end_idx = -1
        
        for i in range(start_idx, len(rest)):
            char = rest[i]
            
            if escape_next:
                escape_next = False
                continue
                
            if char == '\\':
                escape_next = True
                continue
                
            if char == '"':
                in_string = not in_string
                continue
                
            if not in_string:
                if char == '{':
                    brace_count += 1
                elif char == '}':
                    brace_count -= 1
                    if brace_count == 0:
                        end_idx = i
                        break
                        
        if end_idx == -1:
            try:
                return json.loads(rest), ""
            except Exception:
                return {}, text
                
        json_str = rest[start_idx:end_idx+1]
        trailing_text = rest[end_idx+1:].strip()
        
        if trailing_text.startswith('```'):
            trailing_text = trailing_text[3:].strip()
            
        try:
            plot_data = json.loads(json_str)
            return plot_data, trailing_text
        except Exception as e:
            _logger.warning("JSON parse failed in bracket counting: %s", e)
            return {}, text

    def _draw_chart_for_doc(self, data: dict) -> io.BytesIO | None:
        """Creates a premium, physics-aware chart image buffer suitable for injection into DOCX."""
        import numpy as np
        try:
            title  = data.get('title', 'PRC Analysis')
            xAxis_cfg = data.get('xAxis') or {}
            yAxis_cfg = data.get('yAxis') or {}
            yAxis2_cfg = data.get('yAxis2') or {}
            
            xlabel = xAxis_cfg.get('label') or data.get('x_label') or 'X'
            ylabel = yAxis_cfg.get('label') or data.get('y_label') or 'Y'
            ylabel2 = yAxis2_cfg.get('label') or data.get('y_label2') or ''
            
            # 1. Detect if axes are logarithmic
            x_is_log = False
            y_is_log = False
            y2_is_log = False
            
            if data.get('xAxisLog') or xAxis_cfg.get('log'):
                x_is_log = True
            if data.get('yAxisLog') or yAxis_cfg.get('log'):
                y_is_log = True
            if data.get('yAxisRightLog') or yAxis2_cfg.get('log'):
                y2_is_log = True
                
            def check_log_label(lbl):
                if not lbl:
                    return False
                l_str = str(lbl).lower()
                if "relative permeability" in l_str or " kr" in l_str or "kr " in l_str or l_str.startswith("kr"):
                    return False
                keywords = ["permeability", "pc (", "capillary pressure", "resistivity index", "formation factor", "fzi", "rqi"]
                return any(kw in l_str for kw in keywords)
                
            if check_log_label(xlabel):
                x_is_log = True
            if check_log_label(ylabel):
                y_is_log = True
            if ylabel2 and check_log_label(ylabel2):
                y2_is_log = True

            # Extract and build curves
            curves = []
            if 'curves' in data and isinstance(data['curves'], list):
                for c in data['curves']:
                    x_data = c.get('x') or [p.get('x') for p in c.get('data', [])]
                    y_data = c.get('y') or [p.get('y') for p in c.get('data', [])]
                    curves.append({
                        'name': c.get('name') or c.get('label') or 'Series',
                        'x': x_data,
                        'y': y_data,
                        'showLine': c.get('showLine', True),
                        'showPoints': c.get('showPoints', True),
                        'color': c.get('color'),
                        'yId': c.get('yId', 'left')
                    })
            elif 'samples' in data and isinstance(data['samples'], dict):
                for s_name, s_data in data['samples'].items():
                    drainage = s_data.get('drainage', {})
                    x_data = drainage.get('sat_pv') or drainage.get('x', [])
                    y_data = drainage.get('pressure') or drainage.get('y', [])
                    if x_data and y_data:
                        curves.append({
                            'name': s_name,
                            'x': x_data,
                            'y': y_data,
                            'showLine': True,
                            'showPoints': True,
                            'color': None,
                            'yId': 'left'
                        })
            elif 'x' in data and 'y' in data:
                curves.append({
                    'name': data.get('title', 'Series'),
                    'x': data['x'],
                    'y': data['y'],
                    'showLine': True,
                    'showPoints': True,
                    'color': None,
                    'yId': 'left'
                })

            if not curves:
                return None

            # Filter out invalid and non-positive points for log scales
            all_left_x, all_left_y = [], []
            all_right_x, all_right_y = [], []
            
            has_right_axis = data.get('dualAxis', False) or any(c.get('yId') == 'right' for c in curves) or bool(ylabel2)
            
            for curve in curves:
                filtered_x = []
                filtered_y = []
                cur_y_is_log = y2_is_log if (curve['yId'] == 'right' and has_right_axis) else y_is_log
                
                for xi, yi in zip(curve['x'], curve['y']):
                    try:
                        if xi is None or yi is None:
                            continue
                        xf = float(xi)
                        yf = float(yi)
                        if x_is_log and xf <= 0:
                            continue
                        if cur_y_is_log and yf <= 0:
                            continue
                        filtered_x.append(xf)
                        filtered_y.append(yf)
                    except (ValueError, TypeError):
                        continue
                curve['x_filtered'] = filtered_x
                curve['y_filtered'] = filtered_y
                
                if curve['yId'] == 'right' and has_right_axis:
                    all_right_x.extend(filtered_x)
                    all_right_y.extend(filtered_y)
                else:
                    all_left_x.extend(filtered_x)
                    all_left_y.extend(filtered_y)

            # Plot setup
            fig, ax1 = plt.subplots(figsize=(7, 4.5), dpi=150)
            plt.style.use('bmh')
            fig.patch.set_facecolor('white')
            ax1.set_facecolor('#ffffff')
            
            ax2 = None
            if has_right_axis:
                ax2 = ax1.twinx()
                ax2.set_facecolor('none')

            # Color Palette
            _PRC_COLORS = ['#1B3A5C', '#C0392B', '#2E7D32', '#E65100', '#6A1B9A', '#00838F']
            
            # Plot curves
            for i, curve in enumerate(curves):
                x = curve['x_filtered']
                y = curve['y_filtered']
                if not x or not y:
                    continue
                    
                label = curve['name']
                color = curve['color'] or _PRC_COLORS[i % len(_PRC_COLORS)]
                show_line = curve['showLine']
                show_points = curve['showPoints']
                
                target_ax = ax2 if (curve['yId'] == 'right' and ax2) else ax1
                
                # Determine line style and markers
                if not show_line:
                    # Core sample routine - scatter plot with transparency and no line
                    target_ax.scatter(x, y, color=color, label=label, alpha=0.6, s=40, edgecolors='none', zorder=3)
                elif show_line and show_points:
                    target_ax.plot(x, y, marker='o', linestyle='-', linewidth=2.0, markersize=5, color=color, label=label, alpha=0.9, zorder=4)
                else:
                    target_ax.plot(x, y, linestyle='-', linewidth=2.0, color=color, label=label, alpha=0.9, zorder=4)

            # Apply logarithmic scale if needed
            if x_is_log:
                ax1.set_xscale('log')
            if y_is_log:
                ax1.set_yscale('log')
            if y2_is_log and ax2:
                ax2.set_yscale('log')

            # X-axis limits and 5% padding
            is_sat_x = "saturation" in xlabel.lower() or "sw" in xlabel.lower() or "fraction" in xlabel.lower()
            x_min_lim, x_max_lim = None, None
            
            if xAxis_cfg.get('domain') and isinstance(xAxis_cfg['domain'], list) and len(xAxis_cfg['domain']) == 2:
                try:
                    x_min_lim = float(xAxis_cfg['domain'][0])
                    x_max_lim = float(xAxis_cfg['domain'][1])
                except (ValueError, TypeError):
                    pass
                    
            if x_min_lim is None or x_max_lim is None:
                combined_x = all_left_x + all_right_x
                if combined_x:
                    if x_is_log:
                        combined_x = [v for v in combined_x if v > 0]
                        if combined_x:
                            vmin, vmax = min(combined_x), max(combined_x)
                            if vmin == vmax:
                                x_min_lim, x_max_lim = vmin * 0.1, vmax * 10.0
                            else:
                                lmin, lmax = np.log10(vmin), np.log10(vmax)
                                diff = lmax - lmin
                                x_min_lim = 10 ** (lmin - 0.05 * diff)
                                x_max_lim = 10 ** (lmax + 0.05 * diff)
                    else:
                        vmin, vmax = min(combined_x), max(combined_x)
                        if is_sat_x and vmin >= 0 and vmax <= 1.05:
                            x_min_lim, x_max_lim = -0.02, 1.02
                        else:
                            if vmin == vmax:
                                x_min_lim, x_max_lim = vmin - 1.0, vmax + 1.0
                            else:
                                diff = vmax - vmin
                                x_min_lim = vmin - 0.05 * diff
                                x_max_lim = vmax + 0.05 * diff

            if x_min_lim is not None and x_max_lim is not None:
                ax1.set_xlim(x_min_lim, x_max_lim)

            # Primary Y-axis limits
            is_sat_y = "saturation" in ylabel.lower() or "sw" in ylabel.lower() or "fraction" in ylabel.lower()
            y_min_lim, y_max_lim = None, None
            
            if yAxis_cfg.get('domain') and isinstance(yAxis_cfg['domain'], list) and len(yAxis_cfg['domain']) == 2:
                try:
                    y_min_lim = float(yAxis_cfg['domain'][0])
                    y_max_lim = float(yAxis_cfg['domain'][1])
                except (ValueError, TypeError):
                    pass
                    
            if y_min_lim is None or y_max_lim is None:
                if all_left_y:
                    if y_is_log:
                        all_left_y = [v for v in all_left_y if v > 0]
                        if all_left_y:
                            vmin, vmax = min(all_left_y), max(all_left_y)
                            if vmin == vmax:
                                y_min_lim, y_max_lim = vmin * 0.1, vmax * 10.0
                            else:
                                lmin, lmax = np.log10(vmin), np.log10(vmax)
                                diff = lmax - lmin
                                y_min_lim = 10 ** (lmin - 0.05 * diff)
                                y_max_lim = 10 ** (lmax + 0.05 * diff)
                    else:
                        vmin, vmax = min(all_left_y), max(all_left_y)
                        if is_sat_y and vmin >= 0 and vmax <= 1.05:
                            y_min_lim, y_max_lim = -0.02, 1.02
                        else:
                            if vmin == vmax:
                                y_min_lim, y_max_lim = vmin - 1.0, vmax + 1.0
                            else:
                                diff = vmax - vmin
                                y_min_lim = vmin - 0.05 * diff
                                y_max_lim = vmax + 0.05 * diff
                                
            if y_min_lim is not None and y_max_lim is not None:
                ax1.set_ylim(y_min_lim, y_max_lim)

            # Secondary Y-axis limits
            if ax2 and all_right_y:
                y2_min_lim, y2_max_lim = None, None
                if yAxis2_cfg.get('domain') and isinstance(yAxis2_cfg['domain'], list) and len(yAxis2_cfg['domain']) == 2:
                    try:
                        y2_min_lim = float(yAxis2_cfg['domain'][0])
                        y2_max_lim = float(yAxis2_cfg['domain'][1])
                    except (ValueError, TypeError):
                        pass
                if y2_min_lim is None or y2_max_lim is None:
                    if y2_is_log:
                        all_right_y = [v for v in all_right_y if v > 0]
                        if all_right_y:
                            vmin, vmax = min(all_right_y), max(all_right_y)
                            if vmin == vmax:
                                y2_min_lim, y2_max_lim = vmin * 0.1, vmax * 10.0
                            else:
                                lmin, lmax = np.log10(vmin), np.log10(vmax)
                                diff = lmax - lmin
                                y2_min_lim = 10 ** (lmin - 0.05 * diff)
                                y2_max_lim = 10 ** (lmax + 0.05 * diff)
                    else:
                        vmin, vmax = min(all_right_y), max(all_right_y)
                        if vmin == vmax:
                            y2_min_lim, y2_max_lim = vmin - 1.0, vmax + 1.0
                        else:
                            diff = vmax - vmin
                            y2_min_lim = vmin - 0.05 * diff
                            y2_max_lim = vmax + 0.05 * diff
                if y2_min_lim is not None and y2_max_lim is not None:
                    ax2.set_ylim(y2_min_lim, y2_max_lim)

            # Titles and labels
            ax1.set_title(title, fontsize=13, fontweight='bold', color='#1B3A5C', pad=15)
            ax1.set_xlabel(xlabel, fontsize=10, fontweight='bold', color='#333333')
            ax1.set_ylabel(ylabel, fontsize=10, fontweight='bold', color='#333333')
            if ax2 and ylabel2:
                ax2.set_ylabel(ylabel2, fontsize=10, fontweight='bold', color='#333333')

            # Grids and styling
            ax1.grid(True, which='both', linestyle='--', color='#E0E0E0', linewidth=0.5, alpha=0.8)
            if ax2:
                ax2.grid(False)

            # Combine legends
            h1, l1 = ax1.get_legend_handles_labels()
            h2, l2 = ax2.get_legend_handles_labels() if ax2 else ([], [])
            if h1 + h2:
                ax1.legend(h1 + h2, l1 + l2, loc='best', framealpha=0.9, fontsize=8.5, facecolor='white', edgecolor='#CCCCCC')

            # Remove extra spines
            for ax in [ax1, ax2] if ax2 else [ax1]:
                ax.spines['top'].set_visible(False)
                ax.tick_params(axis='both', which='both', labelsize=9)

            buf = io.BytesIO()
            fig.savefig(buf, format='png', bbox_inches='tight', dpi=150)
            plt.close(fig)
            buf.seek(0)
            return buf
        except Exception as e:
            _logger.error("Chart render failed: %s", e)
            return None

    # ─────────────────────────────────────────
    # PUBLIC: Build from pre-fetched JSON content
    # ─────────────────────────────────────────
    def build_from_json(self, raw_content: str, file_type: str,
                        well: str = "PRC Well", engineer: str = "Hviel AI") -> str:
        """
        Called by app.py after Claude returns JSON.
        raw_content: the raw text from Claude (may have ```json fences)
        file_type:   'docx', 'xlsx', 'pptx', or 'pdf'
        Returns: filename string (saved in output_dir)
        """
        # Parse JSON
        text = raw_content.strip().replace('```json', '').replace('```', '').strip()
        try:
            content = json.loads(text)
        except (json.JSONDecodeError, ValueError):
            import re
            m = re.search(r'\{.*\}', text, re.DOTALL)
            try:
                content = json.loads(m.group(0)) if m else {}
            except Exception:
                content = {}

        # Inject well/engineer if not already in content
        if 'author' not in content:
            content['author'] = engineer
        if 'title' not in content:
            content['title'] = f'PRC Report — {well}'

        # --- GLOBAL MARKDOWN TABLE INTERCEPTOR ---
        # Catches LLM hallucinations where tables are dumped as text paragraphs
        if 'tables' not in content:
            content['tables'] = []
            
        for sec in content.get('sections', []):
            new_paragraphs = []
            table_buffer = []
            
            def process_table_buffer():
                if len(table_buffer) >= 3:
                    # Valid markdown table!
                    headers = [c.strip() for c in table_buffer[0].split('|') if c.strip()]
                    rows = []
                    for row_str in table_buffer[2:]: # skip separator
                        cols = [c.strip() for c in row_str.split('|') if c.strip()]
                        if cols: rows.append(cols)
                    if headers and rows:
                        content['tables'].append({
                            'caption': 'Data Table (Auto-Recovered)',
                            'headers': headers,
                            'rows': rows
                        })
                        table_buffer.clear()
                        return
                        
                new_paragraphs.extend(table_buffer)
                table_buffer.clear()

            for para in sec.get('paragraphs', []):
                p_str = str(para).strip()
                
                # If there are any markdown pipes inside the string, we might have part of a table or a full table.
                if '|' in p_str:
                    # Break it down dynamically in case Claude combined lines with newlines
                    lines = [l.strip() for l in p_str.split('\n')]
                    valid_table_lines = [l for l in lines if l.count('|') >= 2]
                    
                    if valid_table_lines:
                        table_buffer.extend(valid_table_lines)
                        # We specifically DO NOT flush here! The next paragraph might be the data rows!
                        continue
                        
                # At this point, the string contains no valid table lines. 
                if p_str == "":
                    # Empty space filler by Claude between rows — ignore it to protect the buffer
                    if not table_buffer:
                        new_paragraphs.append(para)
                    continue
                else:
                    # We hit real text! Time to flush the accumulated table buffer completely.
                    if table_buffer:
                        process_table_buffer()
                    new_paragraphs.append(para)
                    
            if table_buffer:
                process_table_buffer()
                
            sec['paragraphs'] = new_paragraphs

        ts = int(time.time())
        fname = f'PRC_{file_type.upper()}_{ts}'

        if file_type == 'docx':
            return self._build_word(content, fname)
        elif file_type == 'xlsx':
            return self._build_excel(content, fname)
        elif file_type == 'pptx':
            return self._build_pptx(content, fname)
        elif file_type == 'pdf':
            return self._build_pdf(content, fname)
        else:
            raise ValueError(f'Unknown file_type: {file_type}')

    # ─────────────────────────────────────────
    # OPTIONAL: Full end-to-end with Claude API
    # ─────────────────────────────────────────
    def handle_request(self, user_message: str, context: str = '') -> dict:
        """Full pipeline: detect type → ask Claude → build file → return path."""
        if not self.api_key:
            return {'success': False, 'error': 'No API key set'}

        file_type = self._detect_type(user_message)
        if not file_type:
            return {'success': False, 'error': 'Could not detect file type'}

        content = self._ask_claude(user_message, context, file_type)
        if not content:
            return {'success': False, 'error': 'Claude API returned no valid content'}

        ts = int(time.time())
        fname = f'hviel_{file_type}_{ts}'
        try:
            if file_type == 'docx':
                path = self._build_word(content, fname)
            elif file_type == 'xlsx':
                path = self._build_excel(content, fname)
            elif file_type == 'pptx':
                path = self._build_pptx(content, fname)
            elif file_type == 'pdf':
                path = self._build_pdf(content, fname)
            return {'success': True, 'path': path, 'type': file_type}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def _detect_type(self, msg: str) -> str | None:
        import re
        m = msg.lower()
        
        # Check for explicit negation commands regarding files/documents
        negation_patterns = [
            'no word', 'no docx', 'no report', 'no document', 'no file',
            'dont give', "don't give", "don't create", 'dont create',
            'without word', 'without report', 'without document',
            'no excel', 'no xlsx', 'no spreadsheet', 'without excel',
            'no powerpoint', 'no pptx', 'no slides', 'no pdf',
            'answer here', 'answer in chat', 'answer them here', 'answere here',
            'no files', 'no reports', 'no documents', 'not give', 'not create',
            'dont generate', "don't generate", 'not generate', 'without any file',
            'without files', 'stop giving', 'stop generating'
        ]
        if any(pat in m for pat in negation_patterns):
            return None

        # Clean punctuation and whitespace for direct short command check
        cleaned = re.sub(r'[^\w\s]', '', m).strip()
        words = cleaned.split()
        
        # File type keywords - prioritized from specific to general
        pdf_keywords = ['pdf']
        xlsx_keywords = ['excel', 'xlsx', 'spreadsheet']
        pptx_keywords = ['powerpoint', 'pptx', 'presentation', 'slides']
        docx_keywords = ['word', 'docx', 'report', 'document']
        
        all_file_keywords = docx_keywords + xlsx_keywords + pptx_keywords + pdf_keywords
        
        # 1. Direct short command check (e.g. "word", "excel please", "docx", "pdf", etc.)
        # If the user sends a very short message (3 words or fewer) containing a file type keyword,
        # we treat it as an explicit request.
        if len(words) <= 3 and any(w in words for w in all_file_keywords):
            if any(w in words for w in pdf_keywords):
                return 'pdf'
            if any(w in words for w in xlsx_keywords):
                return 'xlsx'
            if any(w in words for w in pptx_keywords):
                return 'pptx'
            if any(w in words for w in docx_keywords):
                return 'docx'
                
        # 2. Action verb check for longer conversational messages
        strong_generation_verbs = [
            'generate', 'create', 'compile', 'export', 'download', 
            'build', 'convert', 'make me', 'generate me', 'create me',
            'build me', 'give me'
        ]
        
        weak_action_verbs = [
            'give', 'write', 'make', 'output', 'provide', 'send', 
            'save', 'extract', 'prepare', 'put in'
        ]
        
        query_indicators = [
            'why', 'what', 'where', 'who', 'when', 'how', 'which', 
            'did you', 'do you', 'is there', 'are there', 'why did', 
            'why is', 'why does', 'what does', 'what is', 'what are', 
            'explain', 'describe', 'show me', 'list', 'check'
        ]
        
        # Check if the user is asking an informational query
        is_query = any(re.search(r'\b' + re.escape(q) + r'\b', m) for q in query_indicators)
        
        # In queries, restrict action triggers only to strong generation verbs
        allowed_verbs = strong_generation_verbs
        if not is_query:
            allowed_verbs = strong_generation_verbs + weak_action_verbs
            
        has_action = False
        for kw in allowed_verbs:
            if ' ' in kw:
                if kw in m:
                    has_action = True
                    break
            else:
                if re.search(r'\b' + re.escape(kw) + r'\b', m):
                    has_action = True
                    break
                    
        # Check target type keywords with word boundaries
        if has_action:
            if any(re.search(r'\b' + re.escape(w) + r'\b', m) for w in pdf_keywords):
                return 'pdf'
            if any(re.search(r'\b' + re.escape(w) + r'\b', m) for w in xlsx_keywords):
                return 'xlsx'
            if any(re.search(r'\b' + re.escape(w) + r'\b', m) for w in pptx_keywords):
                return 'pptx'
            if any(re.search(r'\b' + re.escape(w) + r'\b', m) for w in docx_keywords):
                return 'docx'
                
        return None


    def _ask_claude(self, user_message: str, context: str, file_type: str) -> dict | None:
        """Call Claude API for structured JSON content."""
        import anthropic
        schemas = {
            'docx': '{"title":"...","subtitle":"...","author":"...","sections":[{"heading":"...","level":1,"paragraphs":["..."],"bullets":["..."]}],"tables":[{"caption":"...","headers":["Col1"],"rows":[["val1"]]}]}',
            'xlsx': '{"title":"...","sheets":[{"name":"SheetName","headers":["Col1"],"rows":[["val1"]],"column_widths":[15]}]}',
            'pptx': '{"title":"...","subtitle":"...","slides":[{"title":"SlideTitle","content":"...","bullets":["..."]}]}',
            'pdf':  '{"title":"...","author":"...","sections":[{"heading":"...","paragraphs":["..."],"bullets":["..."]}],"tables":[{"caption":"...","headers":["Col1"],"rows":[["val1"]]}]}',
        }
        system = (
            f"You are Hviel, PRC AI Petrophysical Specialist. Generate a professional {file_type.upper()} file.\n"
            f"Respond with ONLY valid JSON matching this schema:\n{schemas[file_type]}\n"
            "Use real petrophysical data, proper units (mD, %, psi, m TVDSS). Minimum 3 sections/slides."
        )
        try:
            client = anthropic.Anthropic(api_key=self.api_key)
            resp = client.messages.create(
                model='claude-3-5-sonnet-20241022',
                max_tokens=6000,
                system=system,
                messages=[{'role': 'user',
                           'content': f"REQUEST: {user_message}\nCONTEXT:\n{context}" if context
                                      else f"REQUEST: {user_message}"}]
            )
            text = resp.content[0].text.strip().lstrip('`').rstrip('`').strip()
            if text.startswith('json'):
                text = text[4:].strip()
            return json.loads(text)
        except Exception as e:
            _logger.error("Claude API call failed: %s", e)
            return None

    # ════════════════════════════════════════════════════════
    # FILE BUILDERS
    # ════════════════════════════════════════════════════════

    def _build_word(self, content: dict, filename: str) -> str:
        doc = DocxDocument()

        # Global styles
        ns = doc.styles['Normal']
        ns.font.name = 'Arial'
        ns.font.size = Pt(10.5)
        ns.font.color.rgb = GRAY_RGB
        for lvl in range(1, 4):
            hs = doc.styles[f'Heading {lvl}']
            hs.font.name = 'Arial'
            hs.font.color.rgb = NAVY_RGB
        for sec in doc.sections:
            sec.top_margin    = Cm(2.54)
            sec.bottom_margin = Cm(2.54)
            sec.left_margin   = Cm(2.54)
            sec.right_margin  = Cm(2.54)

        # ── Navy banner cover ──
        banner = doc.add_table(rows=1, cols=1)
        banner.alignment = WD_TABLE_ALIGNMENT.CENTER
        bc = banner.rows[0].cells[0]
        self._shade_cell(bc, NAVY)
        bc.height = Cm(2.5)
        bp = bc.paragraphs[0]
        bp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        bp.space_before = Pt(16)
        r = bp.add_run('PETROLEUM RESEARCH CENTER')
        r.bold = True; r.font.size = Pt(20); r.font.color.rgb = WHITE_RGB; r.font.name = 'Arial'
        bp2 = bc.add_paragraph(); bp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = bp2.add_run('Tripoli, Libya')
        r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD); r2.font.name = 'Arial'

        doc.add_paragraph()

        # Logo placeholder
        for lp in ['prc_logo.png', 'prc_logo.jpg']:
            if os.path.exists(lp):
                try:
                    lp2 = doc.add_paragraph(); lp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    lp2.add_run().add_picture(lp, width=Inches(2.0))
                    break
                except Exception:
                    pass

        # Blue divider line
        dp = doc.add_paragraph(); dp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        dp._element.get_or_add_pPr().append(parse_xml(
            f'<w:pBdr {nsdecls("w")}>'
            f'<w:bottom w:val="single" w:sz="12" w:space="8" w:color="{BLUE}"/>'
            f'</w:pBdr>'
        ))
        doc.add_paragraph()

        # Title
        tp = doc.add_paragraph(); tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        tr = tp.add_run(content.get('title', 'PRC Report'))
        tr.bold = True; tr.font.size = Pt(22); tr.font.color.rgb = NAVY_RGB; tr.font.name = 'Arial'

        if content.get('subtitle'):
            sp = doc.add_paragraph(); sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            sr = sp.add_run(content['subtitle'])
            sr.font.size = Pt(13); sr.font.color.rgb = BLUE_RGB; sr.font.name = 'Arial'

        # Info table
        doc.add_paragraph()
        author   = content.get('author', 'Hviel AI')
        date_str = content.get('date', datetime.now().strftime('%d %B %Y'))
        info_rows = [
            ('Prepared By', author),
            ('Date',        date_str),
            ('Classification', 'CONFIDENTIAL \u2014 PRC Internal Use Only'),
        ]
        itbl = doc.add_table(rows=len(info_rows), cols=2)
        itbl.alignment = WD_TABLE_ALIGNMENT.CENTER; itbl.autofit = False
        for row in itbl.rows:
            row.cells[0].width = Inches(2.0); row.cells[1].width = Inches(4.0)
        for i, (lbl, val) in enumerate(info_rows):
            row = itbl.rows[i]; row.height = Cm(0.65)
            c0 = row.cells[0]; c0.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            self._shade_cell(c0, LIGHT_BLUE)
            p0 = c0.paragraphs[0]; p0.space_before = Pt(2); p0.space_after = Pt(2)
            r0 = p0.add_run(lbl); r0.bold = True; r0.font.name = 'Arial'
            r0.font.size = Pt(10); r0.font.color.rgb = NAVY_RGB
            c1 = row.cells[1]; c1.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            p1 = c1.paragraphs[0]; p1.space_before = Pt(2); p1.space_after = Pt(2)
            r1 = p1.add_run(str(val)); r1.font.name = 'Arial'; r1.font.size = Pt(10)
            r1.font.color.rgb = RED_RGB if i == 2 else GRAY_RGB
            for cell in [c0, c1]:
                self._border_cell(cell)

        doc.add_page_break()

        # ── Header & Footer ──
        section = doc.sections[-1]
        section.different_first_page_header_footer = False

        ftr = section.footer; ftr.is_linked_to_previous = False
        fp  = ftr.paragraphs[0] if ftr.paragraphs else ftr.add_paragraph()
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER; fp.clear()
        fp._element.get_or_add_pPr().append(parse_xml(
            f'<w:pBdr {nsdecls("w")}>'
            f'<w:top w:val="single" w:sz="4" w:space="4" w:color="{BLUE}"/>'
            f'</w:pBdr>'
        ))
        fr1 = fp.add_run('Petroleum Research Center \u2014 Confidential')
        fr1.font.name = 'Arial'; fr1.font.size = Pt(8); fr1.font.color.rgb = LIGHT_GRAY_RGB
        fr2 = fp.add_run('    |    Page ')
        fr2.font.name = 'Arial'; fr2.font.size = Pt(8); fr2.font.color.rgb = LIGHT_GRAY_RGB
        self._page_number_field(fp)

        hdr = section.header; hdr.is_linked_to_previous = False
        hp  = hdr.paragraphs[0] if hdr.paragraphs else hdr.add_paragraph()
        hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT; hp.clear()
        hp._element.get_or_add_pPr().append(parse_xml(
            f'<w:pBdr {nsdecls("w")}>'
            f'<w:bottom w:val="single" w:sz="4" w:space="4" w:color="{BLUE}"/>'
            f'</w:pBdr>'
        ))
        hr = hp.add_run(f'PRC SCAL  |  {author}')
        hr.font.name = 'Arial'; hr.font.size = Pt(7.5)
        hr.font.color.rgb = LIGHT_GRAY_RGB; hr.italic = True

        # ── Sections ──
        sec_num = 0
        import re
        for sec in content.get('sections', []):
            level   = int(sec.get('level', 1))
            heading = sec.get('heading', '')
            
            # Clean AI hallucinations (e.g., if it provides '1. Introduction', we strip '1. ' so we don't get '1. 1. Introduction')
            heading = re.sub(r'^\d+[\.\-\s]+', '', heading).strip()
            
            if level == 1:
                sec_num += 1
                heading = f'{sec_num}. {heading}'

            h = doc.add_heading(level=level)
            h.paragraph_format.space_before = Pt(24 if level == 1 else 16)
            h.paragraph_format.space_after = Pt(8)
            hr2 = h.add_run(heading)
            hr2.font.color.rgb = NAVY_RGB; hr2.font.name = 'Arial'; hr2.bold = True
            hr2.font.size = Pt(16 if level == 1 else 13 if level == 2 else 11)

            if level == 1:
                h._element.get_or_add_pPr().append(parse_xml(
                    f'<w:pBdr {nsdecls("w")}>'
                    f'<w:bottom w:val="single" w:sz="6" w:space="4" w:color="{BLUE}"/>'
                    f'</w:pBdr>'
                ))

            for para in sec.get('paragraphs', []):
                para_str = str(para)
                
                # Plot detection — triggers chart injection
                if '__PRC_PLOT__' in para_str:
                    try:
                        # Extract the JSON payload and trailing text
                        chart_data, trailing_text = self._extract_json_payload(para_str)
                        if chart_data:
                            buf = self._draw_chart_for_doc(chart_data)
                            if buf:
                                p = doc.add_paragraph()
                                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                p.add_run().add_picture(buf, width=Inches(6.0))
                                # Add a professional caption
                                cp = doc.add_paragraph()
                                cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                cr = cp.add_run(f"Figure: {chart_data.get('title', 'Petrophysical Analysis')}")
                                cr.italic = True; cr.font.size = Pt(9); cr.font.color.rgb = GRAY_RGB
                                
                                # If there is trailing commentary text, add it as a separate italicized paragraph
                                if trailing_text:
                                    tp = doc.add_paragraph()
                                    tp.paragraph_format.space_after = Pt(12)
                                    tp.paragraph_format.line_spacing = 1.15
                                    tr = tp.add_run(trailing_text)
                                    tr.font.name = 'Arial'; tr.font.size = Pt(10.5); tr.font.color.rgb = GRAY_RGB
                                    tr.italic = True
                                continue # Skip adding the raw JSON as a paragraph
                    except Exception as pe:
                        _logger.warning("Skipping plot injection: %s", pe)

                if '|' in para_str and '---' in para_str and '\n' in para_str.strip():
                    # Hallucination detected! Claude put a markdown table in a paragraph.
                    # Let's dynamically rescue it.
                    lines = [l.strip() for l in para_str.strip().split('\n') if '|' in l]
                    if len(lines) >= 3:
                        headers = [c.strip() for c in lines[0].split('|') if c.strip()]
                        rows = []
                        for l in lines[2:]:
                            cols = [c.strip() for c in l.split('|') if c.strip()]
                            if cols: rows.append(cols)
                        
                        if headers and rows:
                            # Re-use the engine's built-in table logic!
                            content.setdefault('tables', []).append({
                                'caption': 'Data Table (Auto-recovered)',
                                'headers': headers,
                                'rows': rows
                            })
                            continue  # Skip adding this as a paragraph!
                
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(12)
                p.paragraph_format.line_spacing = 1.15
                r = p.add_run(para_str)
                r.font.name = 'Arial'; r.font.size = Pt(10.5); r.font.color.rgb = GRAY_RGB

            for bullet in sec.get('bullets', []):
                bp = doc.add_paragraph(style='List Bullet'); bp.clear()
                br = bp.add_run(str(bullet))
                br.font.name = 'Arial'; br.font.size = Pt(10); br.font.color.rgb = GRAY_RGB

        # ── Tables ──
        for tbl in content.get('tables', []):
            doc.add_paragraph()
            if tbl.get('caption'):
                cp = doc.add_paragraph()
                cr = cp.add_run(tbl['caption'])
                cr.bold = True; cr.italic = True; cr.font.name = 'Arial'
                cr.font.size = Pt(10); cr.font.color.rgb = NAVY_RGB

            headers = tbl.get('headers', [])
            rows    = tbl.get('rows', [])
            if not headers:
                continue
            ncols = len(headers)
            col_w = round(6.0 / ncols, 2)

            table = doc.add_table(rows=1 + len(rows), cols=ncols)
            table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.autofit = False
            for row in table.rows:
                for cell in row.cells:
                    cell.width = Inches(col_w)

            # Header row
            hdr_row = table.rows[0]; hdr_row.height = Cm(0.9)
            for ci, h in enumerate(headers):
                cell = hdr_row.cells[ci]
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                self._shade_cell(cell, NAVY); self._border_cell(cell)
                p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                r = p.add_run(str(h)); r.bold = True
                r.font.size = Pt(9); r.font.color.rgb = WHITE_RGB; r.font.name = 'Arial'

            # Data rows
            for ri, row_data in enumerate(rows):
                row_obj = table.rows[ri + 1]; row_obj.height = Cm(0.75)
                for ci, val in enumerate(row_data):
                    cell = row_obj.cells[ci]
                    cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    self._border_cell(cell)
                    if ri % 2 == 1:
                        self._shade_cell(cell, ZEBRA)
                    p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    r = p.add_run(str(val))
                    r.font.size = Pt(9); r.font.color.rgb = GRAY_RGB; r.font.name = 'Arial'

            doc.add_paragraph()

        # Signature block
        doc.add_paragraph(); doc.add_paragraph()
        sep = doc.add_paragraph()
        sep._element.get_or_add_pPr().append(parse_xml(
            f'<w:pBdr {nsdecls("w")}>'
            f'<w:top w:val="single" w:sz="4" w:space="8" w:color="{BLUE}"/>'
            f'</w:pBdr>'
        ))
        sig = doc.add_table(rows=3, cols=2); sig.alignment = WD_TABLE_ALIGNMENT.CENTER
        for row in sig.rows:
            row.cells[0].width = Inches(3.0); row.cells[1].width = Inches(3.0)
        sig_data = [
            ('Prepared by:', author, 'Hviel \u2014 PRC AI Petrophysical Specialist'),
            ('Reviewed by:', 'PRC Technical Review Board', 'Chief Petrophysicist'),
        ]
        for ci, (role, name, title_) in enumerate(sig_data):
            for ri, lbl in enumerate([role, name, title_]):
                p_ = sig.rows[ri].cells[ci].paragraphs[0]
                r_ = p_.add_run(lbl); r_.font.name = 'Arial'
                if ri == 0:
                    r_.font.size = Pt(9); r_.font.color.rgb = LIGHT_GRAY_RGB
                elif ri == 1:
                    r_.bold = True; r_.font.size = Pt(10); r_.font.color.rgb = NAVY_RGB
                else:
                    r_.font.size = Pt(9); r_.font.color.rgb = GRAY_RGB
        for row in sig.rows:
            for cell in row.cells:
                self._border_none(cell)

        path = os.path.join(self.output_dir, f'{filename}.docx')
        doc.save(path)
        return path

    def _build_excel(self, content: dict, filename: str) -> str:
        wb = Workbook()

        hdr_font  = Font(name='Arial', bold=True, color='FFFFFF', size=11)
        hdr_fill  = PatternFill(start_color='1B3A5C', end_color='1B3A5C', fill_type='solid')
        data_font = Font(name='Arial', size=10, color='333333')
        zfill     = PatternFill(start_color='F2F2F2', end_color='F2F2F2', fill_type='solid')
        brd = Border(
            left  =Side(style='thin', color='CCCCCC'),
            right =Side(style='thin', color='CCCCCC'),
            top   =Side(style='thin', color='CCCCCC'),
            bottom=Side(style='thin', color='CCCCCC'),
        )

        meta = wb.active; meta.title = 'Report Info'
        for i, (k, v) in enumerate([
            ('Title',     content.get('title', '')),
            ('Generated', datetime.now().strftime('%Y-%m-%d %H:%M')),
            ('System',    'Hviel \u2014 PRC AI Petrophysical Specialist'),
        ], 1):
            meta.cell(i, 1, k).font = Font(name='Arial', bold=True, size=11, color='1B3A5C')
            meta.cell(i, 2, v).font = data_font
        meta.column_dimensions['A'].width = 18
        meta.column_dimensions['B'].width = 44

        for sheet in content.get('sheets', []):
            ws = wb.create_sheet(title=str(sheet.get('name', 'Data'))[:31])
            headers = sheet.get('headers', [])
            rows    = sheet.get('rows', [])

            for ci, h in enumerate(headers, 1):
                c = ws.cell(1, ci, str(h))
                c.font = hdr_font; c.fill = hdr_fill; c.border = brd
                c.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)

            for ri, row in enumerate(rows, 2):
                for ci, val in enumerate(row, 1):
                    try:
                        val = float(val)
                    except (ValueError, TypeError):
                        pass
                    c = ws.cell(ri, ci, val)
                    c.font = data_font; c.border = brd
                    c.alignment = Alignment(horizontal='center', vertical='center')
                    if ri % 2 == 0:
                        c.fill = zfill

            widths = sheet.get('column_widths')
            if widths:
                for ci, w in enumerate(widths, 1):
                    ws.column_dimensions[get_column_letter(ci)].width = w
            else:
                for ci, h in enumerate(headers, 1):
                    mx = len(str(h))
                    for row in rows:
                        if ci - 1 < len(row):
                            mx = max(mx, len(str(row[ci - 1])))
                    ws.column_dimensions[get_column_letter(ci)].width = min(mx + 4, 42)

            ws.freeze_panes = 'A2'
            ws.row_dimensions[1].height = 22

        if len(wb.sheetnames) > 1:
            wb.active = 1

        path = os.path.join(self.output_dir, f'{filename}.xlsx')
        wb.save(path); return path

    def _build_pptx(self, content: dict, filename: str) -> str:
        prs = Presentation()
        prs.slide_width  = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        N = PptxRGB(0x1B, 0x3A, 0x5C)
        W = PptxRGB(0xFF, 0xFF, 0xFF)
        G = PptxRGB(0x33, 0x33, 0x33)

        # Title slide — navy background
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = N
        tx = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.5), PptxInches(11.333), PptxInches(2.5))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Hviel Presentation')
        p.font.size = PptxPt(40); p.font.bold = True
        p.font.color.rgb = W; p.font.name = 'Arial'; p.alignment = PP_ALIGN.CENTER
        if content.get('subtitle'):
            p2 = tf.add_paragraph()
            p2.text = content['subtitle']
            p2.font.size = PptxPt(18); p2.font.color.rgb = PptxRGB(0xAA, 0xCC, 0xDD)
            p2.font.name = 'Arial'; p2.alignment = PP_ALIGN.CENTER

        for slide_data in content.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = W

            # Navy header bar
            bar = slide.shapes.add_shape(1,
                PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(1.2))
            bar.fill.solid(); bar.fill.fore_color.rgb = N; bar.line.fill.background()

            tx = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.15), PptxInches(12.3), PptxInches(0.9))
            p = tx.text_frame.paragraphs[0]
            p.text = slide_data.get('title', '')
            p.font.size = PptxPt(26); p.font.bold = True
            p.font.color.rgb = W; p.font.name = 'Arial'

            y = 1.5
            if slide_data.get('content'):
                tx = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(y), PptxInches(12.3), PptxInches(5))
                tf = tx.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = slide_data['content']
                p.font.size = PptxPt(17); p.font.color.rgb = G; p.font.name = 'Arial'
                y += 1.2

            if slide_data.get('bullets'):
                tx = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(y), PptxInches(12.3), PptxInches(5))
                tf = tx.text_frame; tf.word_wrap = True
                for i, b in enumerate(slide_data['bullets']):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f'\u2022  {b}'
                    p.font.size = PptxPt(17); p.font.color.rgb = G
                    p.font.name = 'Arial'; p.space_after = PptxPt(6)

            if slide_data.get('table'):
                t = slide_data['table']
                h_, r_ = t.get('headers', []), t.get('rows', [])
                if h_:
                    shape = slide.shapes.add_table(
                        len(r_) + 1, len(h_),
                        PptxInches(0.5), PptxInches(y),
                        PptxInches(12.3), PptxInches(min(len(r_) * 0.5 + 0.6, 5.5))
                    )
                    tbl = shape.table
                    for ci, htext in enumerate(h_):
                        cell = tbl.cell(0, ci)
                        cell.text = htext
                        cell.fill.solid(); cell.fill.fore_color.rgb = N
                        for para in cell.text_frame.paragraphs:
                            para.font.size = PptxPt(11); para.font.bold = True
                            para.font.color.rgb = W; para.font.name = 'Arial'
                    for ri, row in enumerate(r_):
                        for ci, val in enumerate(row):
                            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
                            for para in cell.text_frame.paragraphs:
                                para.font.size = PptxPt(10); para.font.color.rgb = G
                                para.font.name = 'Arial'

        path = os.path.join(self.output_dir, f'{filename}.pptx')
        prs.save(path); return path

    def _build_pdf(self, content: dict, filename: str) -> str:
        path = os.path.join(self.output_dir, f'{filename}.pdf')
        doc = SimpleDocTemplate(path, pagesize=A4,
            leftMargin=2.5*cm, rightMargin=2.5*cm,
            topMargin=2.5*cm, bottomMargin=2.5*cm)

        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle('PRC_Title', parent=styles['Title'],
            fontName='Helvetica-Bold', fontSize=22, textColor=HexColor('#1B3A5C'),
            alignment=TA_CENTER, spaceAfter=20))
        styles.add(ParagraphStyle('PRC_H1', parent=styles['Heading1'],
            fontName='Helvetica-Bold', fontSize=14, textColor=HexColor('#2E75B6'),
            spaceBefore=18, spaceAfter=8))
        styles.add(ParagraphStyle('PRC_Body', parent=styles['Normal'],
            fontName='Helvetica', fontSize=10, textColor=HexColor('#333333'),
            spaceAfter=6, leading=14))
        styles.add(ParagraphStyle('PRC_Bullet', parent=styles['Normal'],
            fontName='Helvetica', fontSize=10, textColor=HexColor('#333333'),
            leftIndent=20, spaceAfter=3, leading=14))

        elems = []
        elems.append(RLParagraph(content.get('title', 'Report'), styles['PRC_Title']))
        elems.append(Spacer(1, 18))

        for sec in content.get('sections', []):
            elems.append(RLParagraph(sec.get('heading', ''), styles['PRC_H1']))
            for p in sec.get('paragraphs', []):
                elems.append(RLParagraph(str(p), styles['PRC_Body']))
            for b in sec.get('bullets', []):
                elems.append(RLParagraph(f'\u2022  {b}', styles['PRC_Bullet']))

        for tbl in content.get('tables', []):
            elems.append(Spacer(1, 10))
            if tbl.get('caption'):
                elems.append(RLParagraph(tbl['caption'], styles['PRC_H1']))
            data = [tbl['headers']] + tbl['rows']
            t = RLTable(data, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND',  (0, 0), (-1, 0),  HexColor('#1B3A5C')),
                ('TEXTCOLOR',   (0, 0), (-1, 0),  HexColor('#FFFFFF')),
                ('FONTNAME',    (0, 0), (-1, 0),  'Helvetica-Bold'),
                ('FONTSIZE',    (0, 0), (-1, -1), 9),
                ('GRID',        (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
                ('ALIGN',       (0, 0), (-1, -1), 'CENTER'),
                ('TOPPADDING',  (0, 0), (-1, -1), 5),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ] + [('BACKGROUND', (0, i), (-1, i), HexColor('#F2F2F2'))
                 for i in range(2, len(data), 2)]))
            elems.append(t)

        doc.build(elems)
        return path

    # ─────────────────────────────────────────
    # CELL HELPERS
    # ─────────────────────────────────────────
    def _shade_cell(self, cell, color_hex: str):
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>')
        cell._element.get_or_add_tcPr().append(shading)

    def _border_cell(self, cell, color='AAAAAA', sz='4'):
        tc = cell._element; tcPr = tc.get_or_add_tcPr()
        tcB = parse_xml(f'<w:tcBorders {nsdecls("w")}></w:tcBorders>')
        for edge in ['top', 'bottom', 'left', 'right']:
            tcB.append(parse_xml(
                f'<w:{edge} {nsdecls("w")} w:val="single" '
                f'w:sz="{sz}" w:space="0" w:color="{color}"/>'
            ))
        tcPr.append(tcB)

    def _border_none(self, cell):
        tc = cell._element; tcPr = tc.get_or_add_tcPr()
        tcB = parse_xml(f'<w:tcBorders {nsdecls("w")}></w:tcBorders>')
        for edge in ['top', 'bottom', 'left', 'right']:
            tcB.append(parse_xml(
                f'<w:{edge} {nsdecls("w")} w:val="none" w:sz="0" w:space="0" w:color="FFFFFF"/>'
            ))
        tcPr.append(tcB)

    def _page_number_field(self, paragraph):
        for xml in [
            f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>',
            None,
            f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>',
        ]:
            run = paragraph.add_run()
            run.font.name = 'Arial'; run.font.size = Pt(8)
            run.font.color.rgb = LIGHT_GRAY_RGB
            if xml:
                run._element.append(parse_xml(xml))
            else:
                run._element.append(parse_xml(
                    f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>'
                ))


# ════════════════════════════════════════════════════════
# STANDALONE TEST — run this file directly
# ════════════════════════════════════════════════════════
if __name__ == '__main__':
    engine = HvielDocEngine(output_dir='test_generated')

    test_doc = {
        'title': 'SCAL Report \u2014 Well T1-31',
        'subtitle': 'Relative Permeability & Capillary Pressure Study',
        'author': 'Eng. Raouf Elkabir',
        'date': datetime.now().strftime('%d %B %Y'),
        'sections': [
            {'heading': 'Introduction', 'level': 1,
             'paragraphs': [
                 'This report presents SCAL results for Well T1-31 in the Al-Jurf Offshore Field.',
                 'Six core plugs were analysed for relative permeability, capillary pressure, and wettability.'
             ],
             'bullets': ['Unsteady-state relative permeability (USS method)',
                         'Mercury injection capillary pressure (MICP)',
                         'Amott spontaneous imbibition wettability index']},
            {'heading': 'Results & Interpretation', 'level': 1,
             'paragraphs': ['All samples confirm water-wet behaviour. Amott indices range from 0.54 to 0.72.'],
             'bullets': ['Krw @ Sor: 0.072\u20130.142 (strongly water-wet)',
                         'Kro @ Swi: 0.748\u20130.923 (excellent oil mobility)',
                         'Estimated waterflood efficiency: 62\u201375%']},
        ],
        'tables': [
            {'caption': 'Table 1 \u2014 Routine Core Analysis Results',
             'headers': ['Sample ID', 'Depth (m)', 'Porosity (%)', 'Kabs (mD)', 'Swi (%)', 'Amott Index'],
             'rows': [
                 ['T131-01', '2847.3', '19.2', '245.6', '24.8', '0.68'],
                 ['T131-02', '2851.7', '22.1', '412.3', '21.5', '0.72'],
                 ['T131-03', '2856.4', '16.8', '128.9', '29.3', '0.61'],
                 ['T131-04', '2860.9', '20.5', '318.7', '23.1', '0.70'],
                 ['T131-05', '2865.2', '14.3',  '76.4', '33.7', '0.54'],
                 ['T131-06', '2869.8', '21.7', '387.1', '22.0', '0.71'],
             ]},
        ],
    }

    test_xlsx = {
        'title': 'Well T1-31 Core Data',
        'sheets': [{
            'name': 'Core Properties',
            'headers': ['Sample ID', 'Depth (m)', 'Porosity (%)', 'Kabs (mD)', 'Swi (%)', 'Amott'],
            'rows': [
                ['T131-01', '2847.3', '19.2', '245.6', '24.8', '0.68'],
                ['T131-02', '2851.7', '22.1', '412.3', '21.5', '0.72'],
                ['T131-03', '2856.4', '16.8', '128.9', '29.3', '0.61'],
            ]
        }]
    }

    test_pptx = {
        'title': 'SCAL Results \u2014 Well T1-31',
        'subtitle': 'Board Presentation \u2014 PRC Libya',
        'slides': [
            {'title': 'Overview', 'bullets': ['6 core samples from Farwah Formation', 'Water-wet system confirmed', 'Good to excellent reservoir quality']},
            {'title': 'Key Data', 'table': {'headers': ['Sample', 'Porosity', 'Permeability'], 'rows': [['T131-01', '19.2%', '245.6 mD'], ['T131-02', '22.1%', '412.3 mD']]}},
        ]
    }

    print('='*55)
    print('  HVIEL DOC ENGINE \u2014 Self-Test')
    print('='*55)

    print('\n1. Word document...')
    p = engine._build_word(test_doc, 'test_word')
    print(f'   \u2713 {p}  ({os.path.getsize(p)/1024:.1f} KB)')

    print('2. Excel spreadsheet...')
    p = engine._build_excel(test_xlsx, 'test_excel')
    print(f'   \u2713 {p}  ({os.path.getsize(p)/1024:.1f} KB)')

    print('3. PowerPoint presentation...')
    p = engine._build_pptx(test_pptx, 'test_pptx')
    print(f'   \u2713 {p}  ({os.path.getsize(p)/1024:.1f} KB)')

    print('4. PDF...')
    p = engine._build_pdf(test_doc, 'test_pdf')
    print(f'   \u2713 {p}  ({os.path.getsize(p)/1024:.1f} KB)')

    print(f'\n{"="*55}')
    print('  ALL PASSED \u2014 HvielDocEngine is ready!')
    print(f'{"="*55}')
