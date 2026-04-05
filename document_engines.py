import time, json, re, io, os
from datetime import datetime
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt as PtxPt
from pptx.dml.color import RGBColor as PtxRGB

# Word
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# PDF
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table as RLTable, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Excel
from openpyxl import Workbook
from openpyxl.styles import Font as XLFont, PatternFill as XLFill, Alignment as XLAlign, Border as XLBorder, Side as XLSide
from openpyxl.utils import get_column_letter

# ── PRC BRAND COLOURS ──
_NAVY   = DocxRGB(0x1B, 0x3A, 0x5C)   # #1B3A5C — true PRC Navy
_BLUE   = DocxRGB(0x2E, 0x75, 0xB6)   # #2E75B6 — PRC Accent Blue
_WHITE  = DocxRGB(0xFF, 0xFF, 0xFF)
_GRAY   = DocxRGB(0x33, 0x33, 0x33)
_NAVY_HEX  = '1B3A5C'
_ZEBRA_HEX = 'EEF2F8'
_RED    = DocxRGB(0xE3, 0x1E, 0x24)   # PRC Red (cover only)

class DocumentEngines:
    @staticmethod
    def _draw_chart_for_doc(data, format='png'):
        """Creates a chart image buffer suitable for injection into DOC/PPT/PDF."""
        _COLORS = ['#1e3a8a', '#E31E24', '#F59E0B', '#16a34a', '#7c3aed', '#0891b2']
        try:
            fig, ax = plt.subplots(figsize=(7, 4.5), dpi=150)
            plt.style.use('bmh')
            title  = data.get('title', 'PRC Analysis')
            xlabel = data.get('x_label', 'X')
            ylabel = data.get('y_label', 'Y')

            curves = data.get('curves')  # multi-curve schema
            if curves and isinstance(curves, list):
                for i, curve in enumerate(curves):
                    ax.plot(curve.get('x', []), curve.get('y', []),
                            marker='o', linestyle='-', linewidth=2.5, markersize=7,
                            color=_COLORS[i % len(_COLORS)], label=curve.get('label', f'Series {i+1}'))
                ax.legend(fontsize=9, framealpha=0.85)
            else:
                ax.plot(data.get('x', []), data.get('y', []),
                        marker='o', linestyle='-', color='#1e3a8a', linewidth=2.5, markersize=8)

            ax.set_title(title, fontsize=15, fontweight='bold', color='#1e3a8a')
            ax.set_xlabel(xlabel, fontsize=11, fontweight='bold')
            ax.set_ylabel(ylabel, fontsize=11, fontweight='bold')
            ax.grid(True, linestyle='--', alpha=0.6)
            buf = io.BytesIO()
            fig.savefig(buf, format=format, bbox_inches='tight')
            plt.close(fig)
            buf.seek(0)
            return buf
        except:
            return None

    @staticmethod
    def build_pptx(well, engineer, json_str):
        prs = Presentation()
        # Create Title Slide (Layout 0)
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "PRC TECHNICAL EVALUATION"
        title_slide.placeholders[1].text = f"Well: {well.upper()}\nPrepared by: {engineer}\n"
        
        # Override Title Slide text colors
        title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PtxRGB(227, 30, 36) # Red
        title_slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = PtxRGB(45, 53, 142) # Blue
        
        try:
            deck = json.loads(json_str)
            slides_data = deck.get('slides', [])
        except:
            slides_data = [{"title": "Analysis Error", "bullets": ["Failed to parse PPTX JSON scheme."]}]

        for slide_data in slides_data:
            # Bullet Slide (Layout 1)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get('title', 'Analysis Details')
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PtxRGB(45, 53, 142)
            
            tf = slide.placeholders[1].text_frame
            bullets = slide_data.get('bullets', [])
            
            # The first bullet populates the existing first paragraph
            if bullets:
                tf.text = bullets[0]
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
            
            # If chart data is provided
            if 'chart' in slide_data:
                buf = DocumentEngines._draw_chart_for_doc(slide_data['chart'])
                if buf:
                    slide.shapes.add_picture(buf, Inches(1), Inches(2.5), width=Inches(8))

        fname = f"PRC_Presentation_{int(time.time())}.pptx"
        prs.save(fname)
        return fname

    @staticmethod
    def build_pdf(well, engineer, md_text):
        fname = f"PRC_Report_{int(time.time())}.pdf"
        doc = SimpleDocTemplate(fname, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Custom Title Style
        title_style = ParagraphStyle(
            'PRCTitle', parent=styles['Heading1'], fontSize=20, textColor=rl_colors.HexColor('#E31E24'),
            alignment=1, spaceAfter=20
        )
        sub_style = ParagraphStyle(
            'PRCSub', parent=styles['Normal'], fontSize=12, textColor=rl_colors.HexColor('#2D358E'),
            alignment=1, spaceAfter=30
        )
        h2_style = ParagraphStyle(
            'PRCH2', parent=styles['Heading2'], fontSize=16, textColor=rl_colors.HexColor('#2D358E'),
            spaceBefore=15, spaceAfter=10
        )
        body_style = ParagraphStyle('PRCBody', parent=styles['Normal'], fontSize=11, spaceAfter=10)
        
        elements = []
        elements.append(Paragraph("PETROLEUM RESEARCH CENTER", title_style))
        elements.append(Paragraph(f"Well / Project: {well.upper()}", sub_style))
        elements.append(Paragraph(f"Prepared by: {engineer}", sub_style))

        for line in md_text.split('\n'):
            line = line.replace('**', '').strip()
            if not line: continue
            
            if line.startswith('#'):
                clean_text = line.replace('#', '').strip()
                elements.append(Paragraph(clean_text, h2_style))
            elif line.startswith('- ') or line.startswith('* '):
                clean_text = line[2:].strip()
                elements.append(Paragraph(f"• {clean_text}", body_style))
            else:
                elements.append(Paragraph(line, body_style))
                
        doc.build(elements)
        return fname

    @staticmethod
    def build_docx(well, engineer, raw_content):
        """
        Build a PRC-branded Word document from Claude's JSON output.
        Falls back to legacy markdown parsing if JSON cannot be decoded.
        """
        # ── Helper: set cell background ──
        def _bg(cell, hex_col):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), hex_col)
            tcPr.append(shd)

        # ── Try to parse Claude JSON ──
        text = raw_content.strip().replace('```json','').replace('```','').strip()
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            # Fallback: find first { ... } block
            m = re.search(r'\{.*\}', text, re.DOTALL)
            try:
                data = json.loads(m.group(0)) if m else None
            except Exception:
                data = None

        doc = Document()

        # ── Global styles ──
        normal = doc.styles['Normal']
        normal.font.name = 'Arial'
        normal.font.size = DocxPt(11)
        normal.font.color.rgb = _GRAY
        for lvl, size, col in [(1, 16, _NAVY), (2, 13, _BLUE), (3, 11, _BLUE)]:
            hs = doc.styles[f'Heading {lvl}']
            hs.font.name = 'Arial'
            hs.font.size = DocxPt(size)
            hs.font.bold = True
            hs.font.color.rgb = col

        # ── Cover Page ──
        # Logo
        for lp in ['prc_logo.png', 'prc_logo.jpg']:
            if os.path.exists(lp):
                try:
                    lp_para = doc.add_paragraph()
                    lp_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    lp_para.add_run().add_picture(lp, width=DocxInches(2.0))
                    break
                except: pass

        title_para = doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        tr = title_para.add_run('PETROLEUM RESEARCH CENTER')
        tr.bold = True; tr.font.size = DocxPt(22)
        tr.font.name = 'Arial'; tr.font.color.rgb = _NAVY

        doc_title = data.get('title', well) if data else well
        sub_para = doc.add_paragraph()
        sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sub_para.add_run(doc_title.upper())
        sr.bold = True; sr.font.size = DocxPt(15)
        sr.font.name = 'Arial'; sr.font.color.rgb = _RED

        meta_para = doc.add_paragraph()
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        mr = meta_para.add_run(f"Prepared by: {engineer}\u2003\u2003Date: {datetime.now().strftime('%B %d, %Y')}")
        mr.font.size = DocxPt(10); mr.font.name = 'Arial'
        mr.font.color.rgb = DocxRGB(0x88, 0x88, 0x88)
        mr.italic = True
        doc.add_page_break()

        # ── JSON-driven body ──
        if data:
            # Sections
            for sect in data.get('sections', []):
                lvl = min(max(int(sect.get('level', 1)), 1), 3)
                doc.add_heading(sect.get('heading', ''), level=lvl)

                content = (sect.get('content') or '').strip()
                if '__PRC_PLOT__' in content:
                    # Parse and inject chart inline
                    try:
                        plot_json_str = content.split('__PRC_PLOT__', 1)[1].strip()
                        _, end = json.JSONDecoder().raw_decode(plot_json_str)
                        chart_data = json.loads(plot_json_str[:end + plot_json_str.index('}') + 1])
                    except Exception:
                        chart_data = None
                    if chart_data:
                        buf = DocumentEngines._draw_chart_for_doc(chart_data)
                        if buf:
                            p = doc.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            p.add_run().add_picture(buf, width=DocxInches(5.8))
                elif content:
                    doc.add_paragraph(content)

                for bullet in (sect.get('bullets') or []):
                    doc.add_paragraph(str(bullet), style='List Bullet')

            # Tables
            for tbl in data.get('tables', []):
                if tbl.get('title'):
                    doc.add_heading(tbl['title'], level=2)
                headers = tbl.get('headers', [])
                rows    = tbl.get('rows', [])
                if not headers: continue

                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = 'Table Grid'
                table.alignment = WD_TABLE_ALIGNMENT.CENTER

                # Header row
                for ci, hdr in enumerate(headers):
                    cell = table.rows[0].cells[ci]
                    cell.text = ''
                    run = cell.paragraphs[0].add_run(str(hdr))
                    run.bold = True; run.font.color.rgb = _WHITE
                    run.font.name = 'Arial'; run.font.size = DocxPt(10)
                    _bg(cell, _NAVY_HEX)

                # Data rows
                for ri, row_data in enumerate(rows):
                    for ci in range(len(headers)):
                        val = row_data[ci] if ci < len(row_data) else ''
                        cell = table.rows[ri + 1].cells[ci]
                        cell.text = str(val)
                        for run in cell.paragraphs[0].runs:
                            run.font.name = 'Arial'; run.font.size = DocxPt(10)
                        if ri % 2 == 1:
                            _bg(cell, _ZEBRA_HEX)

                doc.add_paragraph()  # spacer

        else:
            # ── Legacy markdown fallback ──
            doc.add_heading('EXECUTIVE SUMMARY & ANALYSIS', 1)
            table_lines = []

            def _render_md_table():
                if not table_lines: return
                parsed = []
                for tl in table_lines:
                    tl = tl.strip().strip('|')
                    parsed.append([c.strip() for c in tl.split('|')])
                actual = [r for i, r in enumerate(parsed) if i != 1]
                if not actual: return
                ncols = max(len(r) for r in actual)
                tbl = doc.add_table(rows=len(actual), cols=ncols)
                tbl.style = 'Table Grid'
                for ri, row in enumerate(actual):
                    for ci in range(ncols):
                        val = row[ci] if ci < len(row) else ''
                        c = tbl.cell(ri, ci)
                        c.text = val
                        if ri == 0:
                            _bg(c, _NAVY_HEX)
                            if c.paragraphs[0].runs:
                                c.paragraphs[0].runs[0].font.bold = True
                                c.paragraphs[0].runs[0].font.color.rgb = _WHITE
                        elif ri % 2 == 0:
                            _bg(c, _ZEBRA_HEX)
                table_lines.clear()
                doc.add_paragraph()

            for line in raw_content.split('\n'):
                line = line.replace('**', '').strip()
                if line.startswith('|') and '|' in line[1:]:
                    table_lines.append(line); continue
                elif table_lines and '|' in line and not line.startswith('#'):
                    table_lines[-1] += ' ' + line; continue
                else:
                    if table_lines: _render_md_table()
                if not line:
                    doc.add_paragraph()
                elif line.startswith('#'):
                    lvl = len(line.split(' ')[0])
                    doc.add_heading(line.replace('#','').strip(), level=min(lvl,3))
                elif line.startswith(('* ','- ')):
                    doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    doc.add_paragraph(line)
            if table_lines: _render_md_table()

        # ── Footer ──
        doc.add_paragraph()
        fp = doc.add_paragraph()
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fr = fp.add_run('This document was generated by Hviel — PRC AI Petrophysical Specialist')
        fr.font.size = DocxPt(8); fr.italtic = True
        fr.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)
        fr.font.name = 'Arial'

        fname = f"PRC_Report_{int(time.time())}.docx"
        doc.save(fname)
        return fname


    @staticmethod
    def build_excel(well, engineer, raw_content):
        """
        Build a PRC-branded Excel workbook from Claude's JSON output.
        Falls back to CSV parsing if JSON cannot be decoded.
        """
        # ── Try to parse Claude JSON ──
        text = raw_content.strip().replace('```json','').replace('```','').strip()
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            m = re.search(r'\{.*\}', text, re.DOTALL)
            try:
                data = json.loads(m.group(0)) if m else None
            except Exception:
                data = None

        # Style constants
        _H_FONT   = XLFont(name='Arial', bold=True, color='FFFFFF', size=11)
        _H_FILL   = XLFill(start_color=_NAVY_HEX, end_color=_NAVY_HEX, fill_type='solid')
        _D_FONT   = XLFont(name='Arial', size=10, color='333333')
        _Z_FILL   = XLFill(start_color='F2F2F2', end_color='F2F2F2', fill_type='solid')
        _BRD      = XLBorder(
            left  =XLSide(style='thin', color='CCCCCC'),
            right =XLSide(style='thin', color='CCCCCC'),
            top   =XLSide(style='thin', color='CCCCCC'),
            bottom=XLSide(style='thin', color='CCCCCC')
        )
        _CA = XLAlign(horizontal='center', vertical='center', wrap_text=True)
        _LA = XLAlign(horizontal='left',   vertical='center', wrap_text=True)

        wb = Workbook()

        # ── Metadata sheet ──
        meta = wb.active
        meta.title = 'Report Info'
        for i, (k, v) in enumerate([
            ('Document Title', data.get('title', well) if data else well),
            ('Well / Project', well),
            ('Prepared By',    engineer),
            ('Date',           datetime.now().strftime('%B %d, %Y')),
            ('System',         'Hviel — PRC AI Petrophysical Specialist'),
        ], 1):
            kc = meta.cell(row=i, column=1, value=k)
            kc.font = XLFont(name='Arial', bold=True, size=11, color=_NAVY_HEX)
            vc = meta.cell(row=i, column=2, value=v)
            vc.font = _D_FONT
        meta.column_dimensions['A'].width = 22
        meta.column_dimensions['B'].width = 44

        # ── JSON-driven sheets ──
        if data and data.get('sheets'):
            for sheet_data in data['sheets']:
                ws = wb.create_sheet(title=str(sheet_data.get('name','Sheet'))[:31])
                headers = sheet_data.get('headers', [])
                rows    = sheet_data.get('rows', [])

                # Header row
                for ci, hdr in enumerate(headers, 1):
                    c = ws.cell(row=1, column=ci, value=str(hdr))
                    c.font = _H_FONT; c.fill = _H_FILL
                    c.border = _BRD; c.alignment = _CA

                # Data rows
                for ri, row_data in enumerate(rows, 2):
                    for ci, val in enumerate(row_data, 1):
                        c = ws.cell(row=ri, column=ci, value=val)
                        c.font = _D_FONT; c.border = _BRD
                        c.alignment = _LA if isinstance(val, str) else _CA
                        if ri % 2 == 0:
                            c.fill = _Z_FILL

                # Auto-width
                for ci, hdr in enumerate(headers, 1):
                    max_len = len(str(hdr))
                    for row_data in rows:
                        if ci - 1 < len(row_data):
                            max_len = max(max_len, len(str(row_data[ci-1])))
                    ws.column_dimensions[get_column_letter(ci)].width = min(max_len + 4, 42)

                ws.freeze_panes = 'A2'
                ws.row_dimensions[1].height = 22

        else:
            # ── CSV fallback ──
            import pandas as pd
            try:
                df = pd.read_csv(io.StringIO(raw_content.strip()))
            except Exception:
                df = pd.DataFrame({'RAW': raw_content.strip().split('\n')})

            ws = wb.create_sheet('PRC SCAL Data')
            for ci, col in enumerate(df.columns, 1):
                c = ws.cell(row=1, column=ci, value=str(col))
                c.font = _H_FONT; c.fill = _H_FILL
                c.border = _BRD; c.alignment = _CA
            for ri, row in enumerate(df.itertuples(index=False), 2):
                for ci, val in enumerate(row, 1):
                    c = ws.cell(row=ri, column=ci, value=val)
                    c.font = _D_FONT; c.border = _BRD; c.alignment = _LA
                    if ri % 2 == 0: c.fill = _Z_FILL
            for col_cells in ws.columns:
                max_len = max((len(str(c.value or '')) for c in col_cells), default=10)
                ws.column_dimensions[col_cells[0].column_letter].width = min(max_len + 4, 45)
            ws.freeze_panes = 'A2'

        fname = f"PRC_Dataset_{int(time.time())}.xlsx"
        wb.save(fname)
        return fname
